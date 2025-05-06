from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from .slide import HTMLSlide, ParagraphContent, BulletNode, BulletTreeContent, TableContent
from .xml_parser import XmlParser 
import os


# Define what placeholder types are considered slide titles
TITLE_PLACEHOLDERS = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}


class SlideConverter:
    """
    Converts a PowerPoint (.pptx) file into Reveal.js-compatible slides.
    Uses both pptx library and XML analysis for reliable content type detection.
    """

    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.slides = []

    def convert(self):
        # Load pptx and extract XML
        prs = Presentation(self.pptx_path)
        try:
            xml_parser = XmlParser(self.pptx_path)
        except Exception as e:
            raise RuntimeError(f"Failed to initialize XML parser: {e}")

        for i, pptx_slide in enumerate(prs.slides):
            # Extract rich paragraph data from XML for slide i
            try:
                shapes = xml_parser.get_slide_shapes(i)
            except Exception as e:
                raise RuntimeError(f"Failed to extract paragraphs from slide {i}: {e}")

            """ 
            # Print structured data for visual debugging
            import json
            print(f"\n✅ Slide {i} Paragraphs Summary:")
            print(json.dumps(shapes, indent=4))
            continue # skip conversion/rendering for now """

            try:
                slide = self.convert_slide(pptx_slide, shapes)
            except Exception as e:
                raise RuntimeError(f"Failed to convert slide {i}: {e}")
            self.slides.append(slide)

    def convert_slide(self, pptx_slide, shapes_data):
        """
        Given a pptx slide and pre-parsed shapes metadata,
        group and convert it into HTML content using Reveal.js structure.
        """

        title = ""
        contents = []

        # Get the title from parsed XML data
        if shapes_data and shapes_data[0].get("title", False):
            title = shapes_data[0]["text"]
            shapes_data = shapes_data[1:] # remove it from body content    

        # Group and convert content with buffer + flush method
        buffer = []
        last_type = None

        def flush_buffer():
            nonlocal buffer, last_type
            if not buffer:
                return
            if last_type in ("bullet", "numbered"):
                ordered = (last_type == "numbered")
                root = BulletNode("ROOT", -1, ordered) # dummy root to hold top level bullets
                stack = [root]
                for item in buffer:
                    node = BulletNode(item["text"], item["level"], ordered)
                    while stack and stack[-1].level >= node.level:
                        stack.pop()
                    stack[-1].add_child(node)
                    stack.append(node)
                contents.append(BulletTreeContent(root))
            elif last_type == "paragraph":
                for item in buffer:
                    contents.append(ParagraphContent(item["text"]))
            buffer.clear()

        # Process all items from XML
        for item in shapes_data:
            # ✅ Handle table content directly
            if item["type"] == "table":
                flush_buffer()  # flush pending bullets/paragraphs first
                contents.append(TableContent(rows=item["rows"],
                                             x_percent=item.get("x_percent"),
                                             y_percent=item.get("y_percent"),
                                             width_percent=item.get("width_percent"),
                                             height_percent=item.get("height_percent"),
                                             col_widths=item.get("col_widths")
                                            ))
                last_type = None  # reset buffer tracking
                continue # skip rest of this loop

            # 🔁 Handle textual content (paragraphs/bullets)
            text = item["text"]
            bullet_type = item["bullet_type"]
            level = item["level"]

            if bullet_type == "number":
                current_type = "numbered"
            elif bullet_type == "bullet":
                current_type = "bullet"
            else:
                current_type = "paragraph"

            if last_type is None or current_type == last_type:
                buffer.append(item)
                last_type = current_type
            else:
                flush_buffer()
                buffer = [item]
                last_type = current_type

        flush_buffer()  # Final flush at end

        slide = HTMLSlide(title=title, transition="fade")
        for content in contents:
            slide.add_content(content)
        return slide

    def save(self, output_file="slides.html"):
        """
        Write all converted slides into a Reveal.js-compatible HTML file.
        """
        try:
            os.makedirs("static", exist_ok=True)
            output_path = os.path.join("static", output_file)
            with open(output_path, "w", encoding="utf-8") as f:
                for slide in self.slides:
                    f.write(slide.to_html())
        except Exception as e:
            raise RuntimeError(f"Failed to write HTML to slides.html: {e}")