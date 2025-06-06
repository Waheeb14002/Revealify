import os
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

#EMU_PER_SLIDE_WIDTH = 9144000
#EMU_PER_SLIDE_HEIGHT = 6858000


class PptxParser:
    def __init__(self, pptx_path):
        self.prs = Presentation(pptx_path)
        self.slide_width = self.prs.slide_width  # in EMUs
        self.slide_height = self.prs.slide_height  # in EMUs

    def get_slide_count(self):
        if self.prs is not None:
            return len(self.prs.slides)
        else: 
            raise RuntimeError("PptxParser Class: No PowerPoint File Has Been Initialized Yet!")

    def _get_bullet_type(self, paragraph):
        p_xml = paragraph._element
        pPr = p_xml.find(qn('a:pPr'))

        if pPr is not None:
            if pPr.find(qn('a:buNone')) is not None:
                return None
            if pPr.find(qn('a:buAutoNum')) is not None:
                return "number"
            if pPr.find(qn('a:buChar')) is not None:
                return "bullet"

        return "bullet"

    def get_slide_shapes(self, slide_index):
        slide = self.prs.slides[slide_index]
        shapes = []

        for shape in slide.shapes:
            x = shape.left
            y = shape.top
            cx = shape.width
            cy = shape.height

            # Detect placeholder type: title, ctrTitle, subTitle, or None
            placeholder_type = None
            if shape.is_placeholder:
                ph_type = getattr(shape.placeholder_format, "type", None)
                if ph_type == PP_PLACEHOLDER.TITLE:
                    placeholder_type = "title"
                elif ph_type == PP_PLACEHOLDER.CENTER_TITLE:
                    placeholder_type = "ctrTitle"
                elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                    placeholder_type = "subTitle"

            # Base shape object with geometry
            shape_obj = {
                "type": None,  # Will be set dynamically
                "title": placeholder_type,
                "x": x,
                "y": y,
                "cx": cx,
                "cy": cy,
                "x_percent": max((x / self.slide_width) * 100, 0),
                "y_percent": max((y / self.slide_height) * 100,0),
                "width_percent": min((cx / self.slide_width) * 100, 100),
                "height_percent": min((cy / self.slide_height) * 100, 100)
            }

            # === Text shapes ===
            if shape.has_text_frame:
                shape_obj["type"] = "text"
                shape_obj["contents"] = []

                for para in shape.text_frame.paragraphs:
                    para_obj = self._parse_paragraph(
                        para,
                        is_title=placeholder_type in ("title", "ctrTitle", "subTitle")
                    )
                    if para_obj:
                        if not para_obj.get("alignment"):
                            para_obj["alignment"] = (
                                "center" if placeholder_type in ("ctrTitle", "subTitle")
                                else "left"
                            )
                        shape_obj["contents"].append(para_obj)

                if shape_obj["contents"]:
                    shapes.append(shape_obj)

            # === Table shapes ===
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                shape_obj["type"] = "table"
                table_data = self._parse_table(shape)
                shape_obj["rows"] = table_data["rows"]
                shape_obj["col_widths"] = table_data["col_widths"]
                shapes.append(shape_obj)

            # === Pictures ===
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                os.makedirs("static/images", exist_ok=True)
                img = shape.image
                ext = img.ext or "png"
                image_bytes = img.blob  # The raw bytes of the image file
                image_name = f"slide{slide_index+1}_img{len(shapes)+1}.{ext}"
                image_path = os.path.join("static/images/", image_name)
                # Save image to static/images/
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                shape_obj["type"] = "image"
                shape_obj["image_path"] = image_path
                shape_obj["image_ext"] = ext
                shape_obj["image_width_px"] = img.size[0]
                shape_obj["image_height_px"] = img.size[1]
                
                shapes.append(shape_obj)

            # === Future: Add chart, image, etc. here ===
            # elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            #     shape_obj["type"] = "chart"
            #     chart_data = self._parse_chart(shape)
            #     shape_obj.update(chart_data)
            #     shapes.append(shape_obj)

        return shapes


    def _parse_paragraph(self, para, is_title=False):
        runs = []
        para_font_size = para.font.size.pt if para.font.size else None
        for run in para.runs:
            run_font_size = run.font.size.pt if run.font.size else None
            font_size_pt = run_font_size if run_font_size else para_font_size
            font_size_px = PptxParser.pt_to_px(font_size_pt) if font_size_pt else None
            run_obj = {
                "text": run.text,
                "bold": run.font.bold,
                "italic": run.font.italic,
                "underline": run.font.underline,
                "hyperlink": run.hyperlink.address if run.hyperlink else None,
                "font_size_px": font_size_px
            }
            runs.append(run_obj)

        if not runs or not any(run["text"].strip() for run in runs):
            return None

        bullet_type = None if is_title else self._get_bullet_type(para)

        align_map = {
            PP_ALIGN.LEFT: "left",
            PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right",
            PP_ALIGN.JUSTIFY: "justify"
            # You can add more if needed
        }
        align = align_map.get(para.alignment, None)

        return {
            "type": "paragraph" if bullet_type is None else "bullet",
            "bullet_type": bullet_type,
            "level": para.level,
            "alignment": align,
            "runs": runs
        }
    

    def _parse_table(self, shape):
        table = shape.table
        rows_data = []
        col_widths = [col.width for col in table.columns]
        total_width = sum(col_widths)
        col_widths_percent = [(w / total_width) * 100 for w in col_widths]

        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_runs = []
                if cell.text_frame:
                    for para in cell.text_frame.paragraphs:
                        para_font_size = para.font.size.pt if para.font.size else None  # fallback
                        for run in para.runs:
                            run_font_size = run.font.size.pt if run.font.size else None
                            font_size_pt = run_font_size if run_font_size else para_font_size
                            font_size_px = PptxParser.pt_to_px(font_size_pt) if font_size_pt else None
                            run_obj = {
                                "text": run.text,
                                "bold": run.font.bold,
                                "italic": run.font.italic,
                                "underline": run.font.underline,
                                "hyperlink": run.hyperlink.address if run.hyperlink else None,
                                "font_size_px": font_size_px
                            }
                            cell_runs.append(run_obj)
                row_data.append(cell_runs)

            rows_data.append(row_data)

        x = shape.left
        y = shape.top
        cx = shape.width
        cy = shape.height

        return {
            "type": "table",
            "rows": rows_data,
            "x": x,
            "y": y,
            "cx": cx,
            "cy": cy,
            "col_widths": col_widths_percent,
            "x_percent": (x / self.slide_width) * 100,
            "y_percent": (y / self.slide_height) * 100,
            "width_percent": (cx / self.slide_width) * 100,
            "height_percent": (cy / self.slide_height) * 100
        }
    
    @staticmethod
    def pt_to_px( pt):
        """
        Convert font size from points (pt) to pixels (px), as used in web browsers.

        PowerPoint and most desktop publishing tools define font sizes in points (pt).
        - 1 point (pt) = 1/72 of an inch.

        Web browsers, by default, use a screen resolution of 96 DPI (dots per inch, or pixels per inch).
        - 1 inch = 96 pixels (px) in CSS.

        To convert points to pixels:
            1. Convert pt to inches:           inches = pt / 72
            2. Convert inches to pixels:       px = inches * 96
            3. Combine into one formula:       px = pt * (96 / 72) = pt * 1.33333...

        Example:
            12 pt * (96 / 72) = 16 px

        Args:
            pt (float or int): The font size in points (pt), as extracted from PowerPoint.

        Returns:
            float: The equivalent font size in pixels (px) for use in web/CSS.
        """
        return pt * 96 / 72  # = pt * 1.33333...



if __name__ == "__main__":
    """
    from pptx.enum import shapes
    print(dir(shapes))
    print("++++++++++++++++++++++++++++++++++++++++++++++++++++++")"""
    import json
    parser = PptxParser("example.pptx")  # adjust path as needed
    for i in range(parser.get_slide_count()):
        print(f"\n============ Slide {i+1} Summary ===================")
        shapes = parser.get_slide_shapes(i)
        print(json.dumps(shapes, indent=4, ensure_ascii=False))
